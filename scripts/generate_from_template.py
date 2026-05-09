#!/usr/bin/env python3
# scripts/generate_from_template.py
"""
从PPTX模板提取图表库，根据数据推荐图表类型，从零创建新PPT。

流程:
  1. 扫描模板PPTX → 提取图表样式（类型、颜色、尺寸）建立图表库
  2. 解析数据文件 → 提取数据集
  3. chart_recommender 推荐图表类型
  4. 从图表库匹配样式 → 用 native_charts 从零创建图表
  5. 自动布局 → 输出全新PPTX

Usage:
    python scripts/generate_from_template.py --data data1.xlsx -o exports/output.pptx
    python scripts/generate_from_template.py --data data2.xlsx -o exports/output.pptx
"""
from __future__ import annotations

import argparse
import datetime
import glob
import json
import re
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from chart_recommender import analyze_dataset, recommend_chart, suggest_layout
from native_charts import (
    add_native_line_chart,
    add_native_bar_chart,
    add_native_pie_chart,
    add_native_horizontal_bar_chart,
    add_native_stacked_bar_chart,
    add_native_area_chart,
    add_native_radar_chart,
    add_native_table,
    ChartSpec,
    SeriesSpec,
    TableSpec,
    px,
)


# ─── 常量 ───

CANVAS_W = 1280
CANVAS_H = 720
MARGIN = 40
TITLE_H = 80
CHART_AREA_Y = TITLE_H + 10
CHART_AREA_H = CANVAS_H - CHART_AREA_Y - MARGIN
CHART_AREA_W = CANVAS_W - MARGIN * 2
GAP = 30

COLOR_PALETTES = [
    # 商务经典
    ["#2B4570", "#4682B4", "#89C2D9", "#F0A500", "#E85D04", "#6A994E", "#386641", "#BC4749"],
    # 科技蓝紫
    ["#3A0CA3", "#4361EE", "#4CC9F0", "#7209B7", "#F72585", "#4895EF", "#560BAD", "#B5179E"],
    # 自然暖色
    ["#A47864", "#B5C4B1", "#D9A066", "#E3C099", "#4A403A", "#8B5E3C", "#6B8F71", "#C9B99A"],
    # 活力对比
    ["#E63946", "#457B9D", "#1D3557", "#A8DADC", "#F1FAEE", "#2A9D8F", "#E9C46A", "#F4A261"],
    # 深沉专业
    ["#264653", "#2A9D8F", "#E9C46A", "#F4A261", "#E76F51", "#606C38", "#283618", "#DDA15E"],
    # 清新明亮
    ["#06D6A0", "#118AB2", "#073B4C", "#FFD166", "#EF476F", "#26547C", "#FF6B6B", "#4ECDC4"],
    # 渐变蓝绿
    ["#005F73", "#0A9396", "#94D2BD", "#EE9B00", "#CA6702", "#BB3E03", "#AE2012", "#9B2226"],
    # 柔和粉彩
    ["#FF6B6B", "#C0D6DF", "#4ECDC4", "#FFE66D", "#95E1D3", "#F38181", "#AA96DA", "#A8D8EA"],
]

DEFAULT_COLORS = COLOR_PALETTES[0]


def _get_palette(slide_index: int, chart_index: int = 0) -> list:
    """根据页码和图表序号轮换配色方案"""
    idx = (slide_index + chart_index) % len(COLOR_PALETTES)
    return COLOR_PALETTES[idx]


# ─── Excel通用数据解析 ───


def parse_xlsx(xlsx_path: Path) -> Tuple[Dict[str, str], List[Dict[str, Any]]]:
    """通用Excel解析，返回 (报告元信息, 数据集列表)

    报告元信息包含:
      - report_title: 报告大标题（第一行）
      - report_summary: 报告摘要/正文描述（标题之后、第一个章节之前的文本）
    """
    wb = load_workbook(str(xlsx_path), read_only=True, data_only=True)
    ws = wb[wb.sheetnames[0]]
    rows = list(ws.iter_rows(values_only=True))
    wb.close()

    if len(rows) <= 1:
        wb = load_workbook(str(xlsx_path), data_only=True)
        ws = wb[wb.sheetnames[0]]
        rows = list(ws.iter_rows(values_only=True))
        wb.close()

    report_title = ""
    summary_lines: List[str] = []
    meta_row_indices: set = set()

    for ri, row in enumerate(rows):
        if row[0] is None:
            continue
        label = str(row[0]).strip()
        if not label:
            continue
        if _is_section_title(label):
            break
        non_null = sum(1 for v in row if v is not None)
        if non_null >= 3:
            break
        if not report_title:
            report_title = label
            meta_row_indices.add(ri)
        else:
            summary_lines.append(label)
            meta_row_indices.add(ri)

    meta = {
        "report_title": report_title,
        "report_summary": "\n".join(summary_lines),
    }

    datasets = []
    current_section = ""
    i = 0
    while i < len(rows):
        row = rows[i]
        if row[0] is None:
            i += 1
            continue

        if i in meta_row_indices:
            i += 1
            continue

        label = str(row[0]).strip()

        if _is_section_title(label):
            current_section = _extract_section_name(label) or label
            i += 1
            continue

        # 趋势表头检测优先于表格检测（修复"月份|品牌1|品牌2"被错误解析为柱状图）
        if _looks_like_trend_header(row):
            ds = _parse_trend_data_direct(rows, i, current_section or "月度趋势")
            if ds:
                ds["section"] = current_section
                datasets.append(ds)
                i += ds.get("_row_count", 1)
                continue

        next_idx = i + 1
        if next_idx < len(rows) and rows[next_idx][0] is not None:
            next_row = rows[next_idx]
            non_null_count = sum(1 for v in row if v is not None)
            if non_null_count >= 3 and isinstance(next_row[1], (int, float)):
                ds = _parse_table_block(rows, i, label)
                if ds:
                    for d in ds:
                        d["section"] = current_section
                    datasets.extend(ds)
                    i += ds[0].get("_skip_rows", 1) + 1
                    continue
            if non_null_count >= 3 and _is_header_row(row):
                ds = _parse_table_block(rows, i, label)
                if ds:
                    for d in ds:
                        d["section"] = current_section
                    datasets.extend(ds)
                    i += ds[0].get("_skip_rows", 1) + 1
                    continue

        if _looks_like_bar_title(label):
            ds = _parse_bar_data(rows, i, label)
            if ds:
                ds["section"] = current_section
                datasets.append(ds)
                i += len(ds.get("categories", [])) + 1
                continue

        if _looks_like_trend_title(label):
            ds = _parse_trend_data(rows, i, label)
            if ds:
                ds["section"] = current_section
                datasets.append(ds)
                i += ds.get("_row_count", 1) + 1
                continue

        i += 1

    return meta, datasets


def _is_section_title(label: str) -> bool:
    """判断是否为章节标题行（如 '1. 2026年...' 或 '10. ...'）"""
    return bool(re.match(r"^\d+\.\s", label))


def _extract_section_name(label: str) -> Optional[str]:
    """从 '1. 2026年...' 格式中提取章节名（去掉编号前缀）"""
    m = re.match(r"^\d+\.\s+(.+)$", label)
    return m.group(1) if m else None


def _is_header_row(row: tuple) -> bool:
    """判断是否为表头行（多个非空字符串列）"""
    str_count = sum(1 for v in row if v is not None and isinstance(v, str) and len(str(v).strip()) > 0)
    return str_count >= 3


def _is_numeric_header(header: str) -> bool:
    """判断表头是否为数值类列名"""
    numeric_keywords = ["销量", "数量", "金额", "增长", "占比", "均价", "评分", "满意度",
                        "渗透率", "差值", "出口", "融资", "排名", "序号"]
    return any(k in header for k in numeric_keywords)


def _detect_col_role(col_name: str, values: list, is_pct: bool) -> str:
    """根据列名和数据特征判断图表角色"""
    pct_keywords = ["占比", "渗透率", "份额", "比例", "比重"]
    rank_keywords = ["排名", "TOP", "top"]
    neg_keywords = ["差值", "变动", "变化"]
    score_keywords = ["评分", "满意度", "得分", "指数", "质量", "服务", "体验"]

    if is_pct or any(k in col_name for k in pct_keywords):
        return "composition"
    if any(k in col_name for k in neg_keywords):
        has_neg = any(v < 0 for v in values if isinstance(v, (int, float)))
        if has_neg:
            return "comparison_negative"
    if any(k in col_name for k in score_keywords):
        return "score"
    return "bar"


def _looks_like_trend_header(row: tuple) -> bool:
    """判断当前行是否为趋势数据的表头行（如 '月份 | 比亚迪 | 特斯拉 | ...'）"""
    if row[0] is None:
        return False
    first = str(row[0]).strip()
    if first not in ("月份", "日期", "时间", "周", "季度"):
        return False
    str_count = sum(1 for v in row if v is not None and isinstance(v, str))
    return str_count >= 3


def _parse_trend_data_direct(rows: list, start: int, title: str) -> Optional[Dict[str, Any]]:
    """解析趋势数据，当前行即为表头行（无单独标题行）"""
    header_row = rows[start]
    series_names = []
    for j in range(1, len(header_row)):
        if header_row[j] is not None:
            series_names.append(str(header_row[j]).strip())

    x_axis = []
    series_data = {name: [] for name in series_names}
    data_start = start + 1
    row_count = 0

    for i in range(data_start, len(rows)):
        row = rows[i]
        if row[0] is None:
            break
        row_count += 1

        date_val = row[0]
        if isinstance(date_val, datetime.datetime):
            x_label = date_val.strftime("%#m/%#d") if sys.platform == "win32" else date_val.strftime("%-m/%-d")
        else:
            x_label = str(date_val)
        x_axis.append(x_label)

        for j, name in enumerate(series_names):
            val = row[j + 1] if j + 1 < len(row) else None
            if isinstance(val, (int, float)):
                series_data[name].append(round(float(val), 2))
            else:
                series_data[name].append(0)

    if not x_axis:
        return None

    series = [{"name": name, "values": series_data[name]} for name in series_names if series_data[name]]

    return {
        "title": title if title != str(rows[start][0]).strip() else "月度趋势",
        "x_axis": x_axis,
        "series": series,
        "chart_role": "trend",
        "_row_count": row_count + 1,
    }


def _looks_like_bar_title(label: str) -> bool:
    keywords = ["对比", "汇总", "排名", "分布", "占比"]
    return any(k in label for k in keywords)


def _looks_like_trend_title(label: str) -> bool:
    keywords = ["走势", "趋势", "月度", "周度", "日度", "变化", "月份"]
    return any(k in label for k in keywords)


def _parse_table_block(rows: list, start: int, first_cell: str) -> Optional[List[Dict[str, Any]]]:
    header_row = rows[start]
    headers = [str(h).strip() if h is not None else "" for h in header_row]

    data_rows = []
    i = start + 1
    while i < len(rows):
        row = rows[i]
        if row[0] is None or str(row[0]).strip() == "":
            break
        data_rows.append(row)
        i += 1

    if not data_rows:
        return None

    category_col = 0
    has_rank_col = headers[0] in ("排名", "序号", "#")
    if has_rank_col and len(headers) > 1:
        for ci in range(1, len(headers)):
            if headers[ci] and not _is_numeric_header(headers[ci]):
                vals = [r[ci] for r in data_rows if ci < len(r) and r[ci] is not None]
                if vals and all(isinstance(v, str) for v in vals):
                    category_col = ci
                    break

    categories = [str(r[category_col]).strip() for r in data_rows]

    numeric_cols = []
    pct_cols = []
    for ci in range(0, len(headers)):
        if ci == category_col:
            continue
        if not headers[ci]:
            continue
        if has_rank_col and ci == 0:
            continue
        vals = [r[ci] for r in data_rows if ci < len(r) and r[ci] is not None]
        if not vals:
            continue
        if not all(isinstance(v, (int, float)) for v in vals):
            continue
        if all(0 <= v <= 1 for v in vals):
            pct_cols.append((ci, headers[ci]))
        elif all(isinstance(v, int) and 1 <= v <= len(data_rows) for v in vals):
            continue
        else:
            numeric_cols.append((ci, headers[ci]))

    is_ranking_table = has_rank_col

    datasets = []

    for ci, col_name in numeric_cols:
        values = [round(float(r[ci]), 2) for r in data_rows if ci < len(r) and isinstance(r[ci], (int, float))]
        if values:
            role = _detect_col_role(col_name, values, False)
            if is_ranking_table and role == "bar":
                role = "ranking"
            datasets.append({
                "title": col_name,
                "categories": categories[:len(values)],
                "values": values,
                "chart_role": role,
                "_skip_rows": len(data_rows),
            })

    for ci, col_name in pct_cols:
        values = [round(float(r[ci]) * 100, 2) for r in data_rows if ci < len(r) and isinstance(r[ci], (int, float))]
        if values:
            role = _detect_col_role(col_name, values, True)
            datasets.append({
                "title": col_name,
                "categories": categories[:len(values)],
                "values": values,
                "unit": "%",
                "chart_role": role,
                "_skip_rows": len(data_rows),
            })

    return datasets if datasets else None


def _parse_bar_data(rows: list, start: int, title: str) -> Optional[Dict[str, Any]]:
    categories = []
    values = []
    i = start + 1
    while i < len(rows):
        row = rows[i]
        if row[0] is None or str(row[0]).strip() == "":
            break
        cat = str(row[0]).strip()
        val = row[1]
        if isinstance(val, (int, float)):
            categories.append(cat)
            values.append(round(float(val), 2))
        i += 1

    if not categories:
        return None

    return {
        "title": title,
        "categories": categories,
        "values": values,
        "chart_role": "bar",
    }


def _parse_trend_data(rows: list, start: int, title: str) -> Optional[Dict[str, Any]]:
    header_idx = start + 1
    if header_idx >= len(rows):
        return None

    header_row = rows[header_idx]
    series_names = []
    for j in range(1, len(header_row)):
        if header_row[j] is not None:
            series_names.append(str(header_row[j]).strip())

    x_axis = []
    series_data = {name: [] for name in series_names}
    data_start = header_idx + 1
    row_count = 0

    for i in range(data_start, len(rows)):
        row = rows[i]
        if row[0] is None:
            break
        row_count += 1

        date_val = row[0]
        if isinstance(date_val, datetime.datetime):
            x_label = date_val.strftime("%#m/%#d") if sys.platform == "win32" else date_val.strftime("%-m/%-d")
        else:
            x_label = str(date_val)
        x_axis.append(x_label)

        for j, name in enumerate(series_names):
            val = row[j + 1] if j + 1 < len(row) else None
            if isinstance(val, (int, float)):
                series_data[name].append(round(float(val), 2))
            else:
                series_data[name].append(0)

    if not x_axis:
        return None

    series = [{"name": name, "values": series_data[name]} for name in series_names if series_data[name]]

    return {
        "title": title,
        "x_axis": x_axis,
        "series": series,
        "chart_role": "trend",
        "_row_count": row_count + 2,
    }


# ─── 图表库（持久化 + 增量扫描） ───

LIBRARY_PATH = SCRIPT_DIR.parent / "chart_library.json"


def _load_library() -> Dict[str, Any]:
    """从磁盘加载持久化图表库"""
    if LIBRARY_PATH.exists():
        try:
            data = json.loads(LIBRARY_PATH.read_text(encoding="utf-8"))
            if "styles" in data and "scanned_files" in data:
                return data
        except (json.JSONDecodeError, KeyError):
            pass
    return {"styles": {}, "scanned_files": {}}


def _save_library(lib_data: Dict[str, Any]) -> None:
    """持久化图表库到磁盘"""
    LIBRARY_PATH.write_text(
        json.dumps(lib_data, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )


def _file_signature(path: Path) -> str:
    """文件签名 = 文件名:大小:修改时间"""
    stat = path.stat()
    return f"{path.name}:{stat.st_size}:{int(stat.st_mtime)}"


def _extract_styles_from_pptx(pptx_path: Path) -> Dict[str, List[Dict[str, Any]]]:
    """从单个PPTX提取图表样式"""
    styles: Dict[str, List[Dict[str, Any]]] = {}
    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        return styles

    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_chart:
                continue

            chart = shape.chart
            ct_str = str(chart.chart_type)
            chart_type = _map_chart_type(ct_str)
            if chart_type == "unknown":
                continue

            series_colors = []
            for plot in chart.plots:
                for series in plot.series:
                    try:
                        fill = series.format.fill
                        if fill.type is not None:
                            rgb = fill.fore_color.rgb
                            series_colors.append(f"#{rgb}")
                    except Exception:
                        pass
                    try:
                        line_color = series.format.line.color.rgb
                        if line_color and f"#{line_color}" not in series_colors:
                            series_colors.append(f"#{line_color}")
                    except Exception:
                        pass

            entry = {
                "source": pptx_path.name,
                "slide": si + 1,
                "shape_name": shape.name,
                "pptx_type": ct_str,
                "series_colors": series_colors if series_colors else None,
                "left": round(shape.left / 9525) if shape.left else 0,
                "top": round(shape.top / 9525) if shape.top else 0,
                "width": round(shape.width / 9525) if shape.width else 400,
                "height": round(shape.height / 9525) if shape.height else 300,
            }

            styles.setdefault(chart_type, []).append(entry)

    return styles


def build_chart_library(template_paths: List[Path]) -> Dict[str, List[Dict[str, Any]]]:
    """增量扫描模板，只处理新增/修改的文件，结果持久化到 chart_library.json"""
    lib_data = _load_library()
    scanned = lib_data["scanned_files"]
    library = lib_data["styles"]

    new_count = 0
    for tpl_path in template_paths:
        if not tpl_path.exists():
            continue
        sig = _file_signature(tpl_path)
        if scanned.get(tpl_path.name) == sig:
            continue

        print(f"  [新模板] 扫描: {tpl_path.name}")
        new_styles = _extract_styles_from_pptx(tpl_path)
        for chart_type, entries in new_styles.items():
            existing = library.setdefault(chart_type, [])
            existing.extend(entries)
        scanned[tpl_path.name] = sig
        new_count += 1

    if new_count > 0:
        _save_library(lib_data)
        print(f"  扫描了 {new_count} 个新模板，已更新图表库")
    else:
        print(f"  图表库已是最新，无需重新扫描")

    return library


def learn_from_output(pptx_path: Path) -> int:
    """从生成的PPT中学习图表样式，扩充图表库"""
    if not pptx_path.exists():
        return 0

    lib_data = _load_library()
    sig = _file_signature(pptx_path)
    if lib_data["scanned_files"].get(pptx_path.name) == sig:
        return 0

    new_styles = _extract_styles_from_pptx(pptx_path)
    added = 0
    for chart_type, entries in new_styles.items():
        existing = lib_data["styles"].setdefault(chart_type, [])
        existing_sigs = {
            f"{e['source']}:{e['slide']}:{e['shape_name']}" for e in existing
        }
        for entry in entries:
            entry_sig = f"{entry['source']}:{entry['slide']}:{entry['shape_name']}"
            if entry_sig not in existing_sigs:
                existing.append(entry)
                added += 1

    lib_data["scanned_files"][pptx_path.name] = sig
    _save_library(lib_data)
    return added


def _map_chart_type(ct_str: str) -> str:
    if "RADAR" in ct_str:
        return "radar_chart"
    if "AREA" in ct_str:
        return "area_chart"
    if "DOUGHNUT" in ct_str:
        return "pie_chart"
    if "PIE" in ct_str:
        return "pie_chart"
    if "LINE" in ct_str:
        return "line_chart"
    if "COLUMN_STACKED" in ct_str:
        return "stacked_bar_chart"
    if "COLUMN" in ct_str:
        return "column_chart"
    if "BAR_STACKED" in ct_str:
        return "stacked_bar_chart"
    if "BAR" in ct_str:
        return "horizontal_bar_chart"
    return "unknown"


def _map_recommended_to_library(rec_type: str) -> str:
    """将推荐的图表类型映射到图表库中可用的类型"""
    line_types = {"line_chart", "dual_axis_line_chart", "stacked_area_chart", "area_chart"}
    bar_types = {"bar_chart", "horizontal_bar_chart", "grouped_bar_chart",
                 "stacked_bar_chart", "waterfall_chart", "butterfly_chart"}
    pie_types = {"donut_chart", "pie_chart"}

    if rec_type in line_types:
        return "line_chart"
    if rec_type in bar_types:
        return "column_chart"
    if rec_type in pie_types:
        return "column_chart"
    return "column_chart"


def _get_style_from_library(library: Dict[str, List[Dict[str, Any]]],
                            chart_type: str) -> Optional[Dict[str, Any]]:
    """从图表库中获取匹配的样式"""
    entries = library.get(chart_type, [])
    if entries:
        return entries[0]
    return None


# ─── 布局计算 ───


def calculate_layout(n_charts: int) -> List[Dict[str, float]]:
    """根据图表数量计算布局位置"""
    if n_charts <= 0:
        return []

    if n_charts == 1:
        return [{"x": MARGIN, "y": CHART_AREA_Y, "w": CHART_AREA_W, "h": CHART_AREA_H}]

    if n_charts == 2:
        w = (CHART_AREA_W - GAP) / 2
        return [
            {"x": MARGIN, "y": CHART_AREA_Y, "w": w, "h": CHART_AREA_H},
            {"x": MARGIN + w + GAP, "y": CHART_AREA_Y, "w": w, "h": CHART_AREA_H},
        ]

    if n_charts == 3:
        w = (CHART_AREA_W - GAP * 2) / 3
        return [
            {"x": MARGIN + i * (w + GAP), "y": CHART_AREA_Y, "w": w, "h": CHART_AREA_H}
            for i in range(3)
        ]

    if n_charts == 4:
        w = (CHART_AREA_W - GAP) / 2
        h = (CHART_AREA_H - GAP) / 2
        slots = []
        for r in range(2):
            for c in range(2):
                slots.append({
                    "x": MARGIN + c * (w + GAP),
                    "y": CHART_AREA_Y + r * (h + GAP),
                    "w": w,
                    "h": h,
                })
        return slots

    cols = 3
    rows_needed = (n_charts + cols - 1) // cols
    w = (CHART_AREA_W - GAP * (cols - 1)) / cols
    h = (CHART_AREA_H - GAP * (rows_needed - 1)) / rows_needed
    slots = []
    for i in range(n_charts):
        r = i // cols
        c = i % cols
        slots.append({
            "x": MARGIN + c * (w + GAP),
            "y": CHART_AREA_Y + r * (h + GAP),
            "w": w,
            "h": h,
        })
    return slots


# ─── 分页逻辑 ───

MAX_CHARTS_PER_SLIDE = 4


def _group_datasets_into_slides(datasets: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """按 section 分组，超过 MAX_CHARTS_PER_SLIDE 的章节拆分为多页"""
    slides: List[Dict[str, Any]] = []
    current_section: Optional[str] = None
    current_group: List[Dict[str, Any]] = []

    for ds in datasets:
        section = ds.get("section", "")
        if section != current_section:
            if current_group:
                _flush_group(slides, current_section or "", current_group)
            current_section = section
            current_group = [ds]
        else:
            current_group.append(ds)

    if current_group:
        _flush_group(slides, current_section or "", current_group)

    return slides


def _flush_group(slides: List[Dict[str, Any]], section: str,
                 group: List[Dict[str, Any]]) -> None:
    """将一个章节的数据集按每页最多 MAX_CHARTS_PER_SLIDE 拆分"""
    total_pages = (len(group) + MAX_CHARTS_PER_SLIDE - 1) // MAX_CHARTS_PER_SLIDE
    for i in range(0, len(group), MAX_CHARTS_PER_SLIDE):
        chunk = group[i:i + MAX_CHARTS_PER_SLIDE]
        page_num = i // MAX_CHARTS_PER_SLIDE + 1
        title = section
        if total_pages > 1:
            title += f" ({page_num}/{total_pages})"
        slides.append({"title": title, "datasets": chunk})


# ─── PPT生成 ───


def _add_cover_slide(slide, title: str, subtitle: str = "") -> None:
    """封面页：居中大标题 + 摘要副标题"""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    title_y = CANVAS_H // 2 - 120 if subtitle else CANVAS_H // 2 - 80
    txBox = slide.shapes.add_textbox(
        px(MARGIN), px(title_y),
        px(CANVAS_W - MARGIN * 2), px(100),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            px(MARGIN + 40), px(title_y + 120),
            px(CANVAS_W - MARGIN * 2 - 80), px(200),
        )
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
        p2.font.name = "Microsoft YaHei"
        p2.alignment = PP_ALIGN.CENTER


def _add_page_title(slide, title: str) -> None:
    """章节标题（每页顶部）"""
    from pptx.dml.color import RGBColor

    txBox = slide.shapes.add_textbox(
        px(MARGIN), px(10),
        px(CANVAS_W - MARGIN * 2), px(TITLE_H - 20),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    p.font.name = "Microsoft YaHei"


def _add_chart_title(slide, title: str, slot: Dict[str, float]) -> None:
    """图表标题（图表上方小字）"""
    from pptx.dml.color import RGBColor

    if not title:
        return
    txBox = slide.shapes.add_textbox(
        px(slot["x"]), px(slot["y"]),
        px(slot["w"]), px(22),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x34, 0x49, 0x5E)
    p.font.name = "Microsoft YaHei"


CHART_TITLE_H = 28

CHART_TYPE_MAP = {
    "composition": "pie",
    "ranking": "horizontal_bar",
    "comparison_negative": "horizontal_bar",
    "trend": "line",
    "bar": "column",
    "multi_series_comparison": "stacked_bar",
    "score": "radar",
}

ALTERNATIVE_TYPE = {
    "column": "horizontal_bar",
    "horizontal_bar": "column",
    "pie": "stacked_bar",
    "line": "area",
    "stacked_bar": "column",
    "area": "line",
    "radar": "horizontal_bar",
}


def _resolve_chart_types(datasets: List[Dict[str, Any]]) -> List[str]:
    """根据 chart_role 决定每个数据集的图表类型，然后执行50%多样性约束"""
    types = []
    for ds in datasets:
        role = ds.get("chart_role", "bar")
        chart_type = CHART_TYPE_MAP.get(role, "column")
        types.append(chart_type)

    if len(types) <= 1:
        return types

    # 50%约束：同一类型不超过半数
    max_same = len(types) // 2 + (1 if len(types) % 2 else 0)
    from collections import Counter
    counts = Counter(types)

    for chart_type, count in counts.items():
        if count > max_same:
            alt = ALTERNATIVE_TYPE.get(chart_type, "horizontal_bar")
            excess = count - max_same
            for i in range(len(types) - 1, -1, -1):
                if excess <= 0:
                    break
                if types[i] == chart_type:
                    types[i] = alt
                    excess -= 1

    return types


def create_pptx(datasets: List[Dict[str, Any]],
                recommendations: List[Dict[str, Any]],
                library: Dict[str, List[Dict[str, Any]]],
                output_path: Path,
                title: str = "",
                subtitle: str = "") -> Tuple[Path, int]:
    """从零创建全新PPTX，多页自动分页，图表类型多样化"""
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = px(CANVAS_W)
    prs.slide_height = px(CANVAS_H)
    blank_layout = prs.slide_layouts[6]

    if title:
        cover = prs.slides.add_slide(blank_layout)
        _add_cover_slide(cover, title, subtitle)

    for i, ds in enumerate(datasets):
        ds["_rec"] = recommendations[i] if i < len(recommendations) else {"primary": {"type": "bar_chart"}}

    slide_specs = _group_datasets_into_slides(datasets)
    total_slides = len(slide_specs)

    for si, slide_spec in enumerate(slide_specs):
        slide = prs.slides.add_slide(blank_layout)
        _add_page_title(slide, slide_spec["title"])

        slide_datasets = slide_spec["datasets"]
        n = len(slide_datasets)
        slots = calculate_layout(n)

        chart_types = _resolve_chart_types(slide_datasets)

        for j, ds in enumerate(slide_datasets):
            if j >= len(slots):
                break

            slot = slots[j]
            ds.pop("_rec", None)

            colors = _get_palette(si, j)

            chart_title = ds.get("title", "")
            _add_chart_title(slide, chart_title, slot)

            chart_slot = {
                "x": slot["x"],
                "y": slot["y"] + CHART_TITLE_H,
                "w": slot["w"],
                "h": slot["h"] - CHART_TITLE_H,
            }

            resolved_type = chart_types[j]

            if resolved_type == "line" or ds.get("chart_role") == "trend":
                _add_trend_chart(slide, ds, chart_slot, colors, j)
            elif resolved_type == "area":
                _add_area_chart(slide, ds, chart_slot, colors, j)
            elif resolved_type == "pie":
                _add_pie_chart(slide, ds, chart_slot, colors, j)
            elif resolved_type == "horizontal_bar":
                _add_horizontal_bar_chart(slide, ds, chart_slot, colors, j)
            elif resolved_type == "stacked_bar":
                _add_stacked_bar_chart(slide, ds, chart_slot, colors, j)
            elif resolved_type == "radar":
                _add_radar_chart(slide, ds, chart_slot, colors, j)
            else:
                _add_bar_chart(slide, ds, chart_slot, colors, j, "column_chart")

        type_summary = ", ".join(f"{t}" for t in dict.fromkeys(chart_types))
        print(f"  [Slide {si + 1}] {slide_spec['title']} — {n} 个图表 ({type_summary})")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path, total_slides


def _add_trend_chart(slide, dataset: Dict[str, Any], slot: Dict[str, float],
                     colors: List[str], index: int) -> None:
    """添加趋势折线图"""
    series_list = dataset.get("series", [])
    categories = dataset.get("x_axis", [])

    spec_series = []
    for j, s in enumerate(series_list):
        color = colors[j % len(colors)]
        spec_series.append(SeriesSpec(
            name=s["name"],
            values=s["values"],
            color=color,
        ))

    all_vals = [v for s in series_list for v in s["values"] if isinstance(v, (int, float))]
    if all_vals:
        v_min = min(all_vals)
        v_max = max(all_vals)
        axis_min = max(0, v_min - (v_max - v_min) * 0.1)
        axis_max = v_max + (v_max - v_min) * 0.1
        rng = axis_max - axis_min
        if rng > 0:
            major_unit = round(rng / 5, -len(str(int(rng // 5))) + 1) if rng > 10 else rng / 5
        else:
            major_unit = 1
    else:
        axis_min, axis_max, major_unit = 0, 100, 20

    spec = ChartSpec(
        slide_index=0,
        native_type="line_chart",
        x=slot["x"],
        y=slot["y"],
        width=slot["w"],
        height=slot["h"],
        categories=categories,
        series=spec_series,
        axis_min=axis_min,
        axis_max=axis_max,
        axis_major_unit=major_unit,
        category_font_size=8,
        value_font_size=9,
        legend_font_size=8,
    )

    add_native_line_chart(slide, spec)
    print(f"  [OK] 创建折线图: {dataset['title']} ({len(series_list)}个系列, {len(categories)}个数据点)")


def _add_bar_chart(slide, dataset: Dict[str, Any], slot: Dict[str, float],
                   colors: List[str], index: int, rec_type: str) -> None:
    """添加柱状图"""
    categories = dataset.get("categories", [])
    values = dataset.get("values", [])

    spec_series = [SeriesSpec(
        name=dataset.get("title", "数值"),
        values=values,
        color=colors[index % len(colors)],
    )]

    spec = ChartSpec(
        slide_index=0,
        native_type="column_chart",
        x=slot["x"],
        y=slot["y"],
        width=slot["w"],
        height=slot["h"],
        categories=categories,
        series=spec_series,
        category_font_size=8,
        value_font_size=9,
        legend_font_size=8,
    )

    add_native_bar_chart(slide, spec)
    print(f"  [OK] 创建柱状图: {dataset['title']} ({len(categories)}个类别)")


def _add_pie_chart(slide, dataset: Dict[str, Any], slot: Dict[str, float],
                   colors: List[str], index: int) -> None:
    """添加环形图"""
    categories = dataset.get("categories", [])
    values = dataset.get("values", [])

    spec_series = [SeriesSpec(
        name=dataset.get("title", "数值"),
        values=values,
        color=json.dumps(colors),
    )]

    spec = ChartSpec(
        slide_index=0,
        native_type="pie_chart",
        x=slot["x"],
        y=slot["y"],
        width=slot["w"],
        height=slot["h"],
        categories=categories,
        series=spec_series,
        legend_font_size=8,
    )

    add_native_pie_chart(slide, spec)
    print(f"  [OK] 创建环形图: {dataset['title']} ({len(categories)}个类别)")


def _add_horizontal_bar_chart(slide, dataset: Dict[str, Any], slot: Dict[str, float],
                              colors: List[str], index: int) -> None:
    """添加横向条形图"""
    categories = dataset.get("categories", [])
    values = dataset.get("values", [])

    spec_series = [SeriesSpec(
        name=dataset.get("title", "数值"),
        values=values,
        color=colors[index % len(colors)],
    )]

    spec = ChartSpec(
        slide_index=0,
        native_type="horizontal_bar_chart",
        x=slot["x"],
        y=slot["y"],
        width=slot["w"],
        height=slot["h"],
        categories=categories,
        series=spec_series,
        category_font_size=8,
        value_font_size=9,
        legend_font_size=8,
    )

    add_native_horizontal_bar_chart(slide, spec)
    print(f"  [OK] 创建横向条形图: {dataset['title']} ({len(categories)}个类别)")


def _add_stacked_bar_chart(slide, dataset: Dict[str, Any], slot: Dict[str, float],
                           colors: List[str], index: int) -> None:
    """添加堆叠柱状图（用占比数据生成多系列堆叠）"""
    categories = dataset.get("categories", [])
    values = dataset.get("values", [])

    total = sum(v for v in values if isinstance(v, (int, float)))
    if total > 0:
        rest_values = [round(total - v, 2) for v in values]
        spec_series = [
            SeriesSpec(name=dataset.get("title", "数值"), values=values, color=colors[0]),
            SeriesSpec(name="其他", values=rest_values, color=colors[2 % len(colors)]),
        ]
    else:
        spec_series = [SeriesSpec(name=dataset.get("title", "数值"), values=values, color=colors[0])]

    spec = ChartSpec(
        slide_index=0,
        native_type="stacked_bar_chart",
        x=slot["x"],
        y=slot["y"],
        width=slot["w"],
        height=slot["h"],
        categories=categories,
        series=spec_series,
        category_font_size=8,
        value_font_size=9,
        legend_font_size=8,
    )

    add_native_stacked_bar_chart(slide, spec)
    print(f"  [OK] 创建堆叠柱状图: {dataset['title']} ({len(categories)}个类别)")


def _add_area_chart(slide, dataset: Dict[str, Any], slot: Dict[str, float],
                    colors: List[str], index: int) -> None:
    """添加面积图"""
    series_list = dataset.get("series", [])
    categories = dataset.get("x_axis", dataset.get("categories", []))

    if series_list:
        spec_series = [
            SeriesSpec(name=s["name"], values=s["values"], color=colors[j % len(colors)])
            for j, s in enumerate(series_list)
        ]
    else:
        values = dataset.get("values", [])
        spec_series = [SeriesSpec(name=dataset.get("title", "数值"), values=values, color=colors[0])]

    all_vals = []
    for s in spec_series:
        all_vals.extend(v for v in s.values if isinstance(v, (int, float)))
    if all_vals:
        v_min = min(all_vals)
        v_max = max(all_vals)
        axis_min = max(0, v_min - (v_max - v_min) * 0.1)
        axis_max = v_max + (v_max - v_min) * 0.1
    else:
        axis_min, axis_max = 0, 100

    spec = ChartSpec(
        slide_index=0,
        native_type="area_chart",
        x=slot["x"],
        y=slot["y"],
        width=slot["w"],
        height=slot["h"],
        categories=categories,
        series=spec_series,
        axis_min=axis_min,
        axis_max=axis_max,
        category_font_size=8,
        value_font_size=9,
        legend_font_size=8,
    )

    add_native_area_chart(slide, spec)
    n_series = len(spec_series)
    print(f"  [OK] 创建面积图: {dataset['title']} ({n_series}个系列, {len(categories)}个数据点)")


def _add_radar_chart(slide, dataset: Dict[str, Any], slot: Dict[str, float],
                     colors: List[str], index: int) -> None:
    """添加雷达图"""
    categories = dataset.get("categories", [])
    values = dataset.get("values", [])

    spec_series = [SeriesSpec(
        name=dataset.get("title", "数值"),
        values=values,
        color=colors[index % len(colors)],
    )]

    spec = ChartSpec(
        slide_index=0,
        native_type="radar_chart",
        x=slot["x"],
        y=slot["y"],
        width=slot["w"],
        height=slot["h"],
        categories=categories,
        series=spec_series,
        legend_font_size=8,
    )

    add_native_radar_chart(slide, spec)
    print(f"  [OK] 创建雷达图: {dataset['title']} ({len(categories)}个维度)")


# ─── 主流程 ───


def main() -> int:
    parser = argparse.ArgumentParser(
        description="从PPTX模板提取图表库 + 数据推荐 + 生成全新PPT",
    )
    parser.add_argument("--data", type=Path, required=True,
                        help="数据文件路径（xlsx）")
    parser.add_argument("--templates", type=str, default=None,
                        help="模板文件通配符（如 'ppt模板*.pptx'）")
    parser.add_argument("-o", "--output", type=Path, default=None,
                        help="输出PPTX路径")
    args = parser.parse_args()

    data_path = args.data.expanduser().resolve()
    if not data_path.exists():
        print(f"[ERROR] 数据文件不存在: {data_path}")
        return 1

    # Step 1: 建立图表库（增量扫描，只处理新模板）
    print(f"\n[Step 1] 加载图表库")
    if args.templates:
        base_dir = Path.cwd()
        template_paths = [Path(p) for p in sorted(glob.glob(str(base_dir / args.templates)))]
    else:
        base_dir = Path.cwd()
        template_paths = sorted(base_dir.glob("ppt模板*.pptx"))

    if not template_paths:
        lib_data = _load_library()
        library = lib_data["styles"]
        if library:
            total = sum(len(v) for v in library.values())
            print(f"  从缓存加载 {total} 个图表样式")
        else:
            print("  [WARN] 未找到模板文件，使用默认样式")
    else:
        library = build_chart_library(template_paths)

    total_styles = sum(len(v) for v in library.values())
    for chart_type, entries in library.items():
        sources = set(e["source"] for e in entries)
        print(f"    {chart_type}: {len(entries)} 个 (来自 {', '.join(sources)})")

    # Step 2: 解析数据
    print(f"\n[Step 2] 解析数据: {data_path.name}")
    meta, datasets = parse_xlsx(data_path)
    if meta["report_title"]:
        print(f"  报告标题: {meta['report_title']}")
    if meta["report_summary"]:
        print(f"  报告摘要: {meta['report_summary'][:80]}...")
    for ds in datasets:
        ds.pop("_row_count", None)
        ds.pop("_skip_rows", None)
    print(f"  提取了 {len(datasets)} 个数据集:")
    sections_seen = {}
    for ds in datasets:
        title = ds.get("title", "未知")
        role = ds.get("chart_role", "unknown")
        section = ds.get("section", "")
        if section and section not in sections_seen:
            sections_seen[section] = True
            print(f"  [章节] {section}")
        if role == "bar":
            print(f"    - [{role}] {title}: {len(ds.get('categories', []))} 个类别")
        else:
            n_series = len(ds.get("series", []))
            n_points = len(ds.get("x_axis", []))
            print(f"    - [{role}] {title}: {n_series} 个系列, {n_points} 个数据点")

    if not datasets:
        print("[ERROR] 未从数据文件中提取到有效数据集")
        return 1

    # Step 3: 图表推荐
    print(f"\n[Step 3] 数据分析与图表推荐")
    recommendations = []
    for ds in datasets:
        features = analyze_dataset(ds)
        rec = recommend_chart(features)
        recommendations.append(rec)
        lib_type = _map_recommended_to_library(rec["primary"]["type"])
        has_style = "有" if _get_style_from_library(library, lib_type) else "无"
        print(f"  {ds['title']} → {rec['primary']['type']} → 库匹配: {lib_type} ({has_style}样式)")

    # Step 4: 创建多页PPT
    print(f"\n[Step 4] 创建多页PPT（自动分页，每页最多 {MAX_CHARTS_PER_SLIDE} 个图表）")

    if args.output:
        output_path = args.output.expanduser().resolve()
    else:
        stem = data_path.stem
        output_path = Path.cwd() / "exports" / f"{stem}_分析报告_{datetime.datetime.now().strftime('%Y%m%d')}.pptx"

    report_title = meta["report_title"] or f"{data_path.stem} 数据分析报告"
    report_summary = meta["report_summary"]
    result, n_slides = create_pptx(datasets, recommendations, library, output_path,
                                   title=report_title, subtitle=report_summary)

    print(f"\n[OK] 生成完成!")
    print(f"  输入数据: {data_path.name}")
    print(f"  图表库: {total_styles} 个样式")
    print(f"  输出文件: {result}")
    print(f"  共 {n_slides + 1} 页（含封面），{len(datasets)} 个原生可编辑图表")

    # Step 5: 从生成的PPT学习图表样式
    learned = learn_from_output(result)
    if learned > 0:
        print(f"  [学习] 从输出PPT中学到 {learned} 个新图表样式，已更新图表库")

    return 0


if __name__ == "__main__":
    sys.exit(main())
