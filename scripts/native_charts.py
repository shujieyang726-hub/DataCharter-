#!/usr/bin/env python3
"""native_charts.py — Post-process PPTX to replace SVG-based static
charts and tables with native PowerPoint objects.

Usage:
    python3 scripts/native_charts.py <project_path> [--pptx <path>]

Reads SVG files from svg_final/ (or svg_output/), looks for <g> elements
with ``data-native-type`` attributes, then opens the exported PPTX and
replaces the corresponding static shape groups with native editable
PowerPoint charts or tables.

Supported data-native-type values:
    line_chart      → native line chart with markers
    bar_chart       → native clustered bar chart (horizontal)
    column_chart    → native clustered column chart (vertical)
    table           → native PowerPoint table

SVG metadata convention (added by Executor on <g> elements):

    <g id="my-chart"
       data-native-type="line_chart"
       data-native-x="40" data-native-y="90"
       data-native-width="1200" data-native-height="320"
       data-native-categories="4/5,4/6,4/7"
       data-native-series='[{"name":"Vol","values":[100,200,300],"color":"#C00000"}]'>
      ... fallback SVG shapes ...
    </g>

    <g id="my-table"
       data-native-type="table"
       data-native-x="60" data-native-y="110"
       data-native-width="1160" data-native-height="500"
       data-native-headers='["Col1","Col2"]'
       data-native-rows='[["a","b"],["c","d"]]'
       data-native-header-color="#C00000"
       data-native-row-colors='["#CFD5EA","#E9EBF5"]'
       data-native-col-widths='[80,200]'>
      ... fallback SVG shapes ...
    </g>
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import sys
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_LABEL_POSITION
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

SVG_NS = "http://www.w3.org/2000/svg"

# 1 SVG px = 9525 EMU  (slide 1280px = 12192000 EMU)
PX_TO_EMU = 9525


def px(val: float) -> int:
    return int(val * PX_TO_EMU)


# ── Data classes ─────────────────────────────────────────────────────────────

@dataclass
class SeriesSpec:
    name: str
    values: list[float]
    color: str | None = None


@dataclass
class ChartSpec:
    slide_index: int
    native_type: str
    x: float
    y: float
    width: float
    height: float
    categories: list[str] = field(default_factory=list)
    series: list[SeriesSpec] = field(default_factory=list)
    axis_min: float | None = None
    axis_max: float | None = None
    axis_major_unit: float | None = None
    value_number_format: str | None = None
    category_font_size: int | None = None
    value_font_size: int | None = None
    legend_font_size: int | None = None
    group_id: str = ""


@dataclass
class TableSpec:
    slide_index: int
    native_type: str = "table"
    x: float = 0
    y: float = 0
    width: float = 0
    height: float = 0
    headers: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)
    header_color: str = "#C00000"
    header_text_color: str = "#FFFFFF"
    body_text_color: str = "#000000"
    row_colors: list[str] = field(default_factory=lambda: ["#CFD5EA", "#E9EBF5"])
    col_widths: list[float] | None = None
    col_alignments: list[str] | None = None
    font_size: int = 10
    header_font_size: int = 10
    group_id: str = ""


# ── SVG parsing ──────────────────────────────────────────────────────────────

def _parse_bounds_from_children(g_elem: ET.Element) -> tuple[float, float, float, float]:
    """Estimate bounding box of a <g> from its child elements' coordinates."""
    min_x = min_y = float("inf")
    max_x = max_y = float("-inf")

    for child in g_elem.iter():
        tag = child.tag.replace(f"{{{SVG_NS}}}", "")
        # rect
        if tag == "rect":
            cx = float(child.get("x", 0))
            cy = float(child.get("y", 0))
            cw = float(child.get("width", 0))
            ch = float(child.get("height", 0))
            min_x = min(min_x, cx)
            min_y = min(min_y, cy)
            max_x = max(max_x, cx + cw)
            max_y = max(max_y, cy + ch)
        # text
        elif tag == "text":
            tx = float(child.get("x", "0").split()[0])
            ty = float(child.get("y", "0").split()[0])
            min_x = min(min_x, tx)
            min_y = min(min_y, ty - 14)  # approximate ascender
            max_x = max(max_x, tx + 100)
            max_y = max(max_y, ty + 4)
        # polyline / polygon
        elif tag in ("polyline", "polygon"):
            pts = child.get("points", "")
            for pair in pts.split():
                parts = pair.split(",")
                if len(parts) == 2:
                    px_val, py_val = float(parts[0]), float(parts[1])
                    min_x = min(min_x, px_val)
                    min_y = min(min_y, py_val)
                    max_x = max(max_x, px_val)
                    max_y = max(max_y, py_val)
        # line
        elif tag == "line":
            for attr in ("x1", "x2"):
                v = child.get(attr)
                if v:
                    min_x = min(min_x, float(v))
                    max_x = max(max_x, float(v))
            for attr in ("y1", "y2"):
                v = child.get(attr)
                if v:
                    min_y = min(min_y, float(v))
                    max_y = max(max_y, float(v))
        # circle
        elif tag == "circle":
            ccx = float(child.get("cx", 0))
            ccy = float(child.get("cy", 0))
            cr = float(child.get("r", 0))
            min_x = min(min_x, ccx - cr)
            min_y = min(min_y, ccy - cr)
            max_x = max(max_x, ccx + cr)
            max_y = max(max_y, ccy + cr)

    if min_x == float("inf"):
        return 0, 0, 100, 100
    return min_x, min_y, max_x - min_x, max_y - min_y


def extract_native_specs(svg_path: Path, slide_index: int) -> list[ChartSpec | TableSpec]:
    """Parse an SVG file and extract all <g data-native-type="..."> specs."""
    tree = ET.parse(str(svg_path))
    root = tree.getroot()

    specs: list[ChartSpec | TableSpec] = []
    seen_ids: set[str] = set()

    # Search both namespaced and non-namespaced <g> elements
    g_elements = list(root.iter(f"{{{SVG_NS}}}g")) + list(root.iter("g"))

    for g in g_elements:
        native_type = g.get("data-native-type")
        if not native_type:
            continue

        group_id = g.get("id", "")
        # Deduplicate (same element may appear in both iter passes)
        dedup_key = f"{group_id}:{id(g)}"
        if dedup_key in seen_ids:
            continue
        seen_ids.add(dedup_key)

        # Position: explicit data-native-x/y/width/height or infer from children
        if g.get("data-native-x") is not None:
            bx = float(g.get("data-native-x", 0))
            by = float(g.get("data-native-y", 0))
            bw = float(g.get("data-native-width", 400))
            bh = float(g.get("data-native-height", 300))
        else:
            bx, by, bw, bh = _parse_bounds_from_children(g)

        if native_type == "table":
            headers_raw = g.get("data-native-headers", "[]")
            rows_raw = g.get("data-native-rows", "[]")
            header_color = g.get("data-native-header-color", "#C00000")
            header_text_color = g.get("data-native-header-text-color", "#FFFFFF")
            body_text_color = g.get("data-native-body-text-color", "#000000")
            row_colors_raw = g.get("data-native-row-colors", '["#CFD5EA","#E9EBF5"]')
            col_widths_raw = g.get("data-native-col-widths")
            col_alignments_raw = g.get("data-native-col-alignments")
            font_size_raw = g.get("data-native-font-size")
            header_font_size_raw = g.get("data-native-header-font-size")

            try:
                headers = json.loads(headers_raw)
                rows = json.loads(rows_raw)
                row_colors = json.loads(row_colors_raw)
                col_widths = json.loads(col_widths_raw) if col_widths_raw else None
                col_alignments = json.loads(col_alignments_raw) if col_alignments_raw else None
            except json.JSONDecodeError as e:
                print(f"  Warning: JSON parse error in {svg_path.name} group '{group_id}': {e}")
                continue

            specs.append(TableSpec(
                slide_index=slide_index,
                x=bx, y=by, width=bw, height=bh,
                headers=headers, rows=rows,
                header_color=header_color,
                header_text_color=header_text_color,
                body_text_color=body_text_color,
                row_colors=row_colors,
                col_widths=col_widths,
                col_alignments=col_alignments,
                font_size=int(font_size_raw) if font_size_raw else 10,
                header_font_size=int(header_font_size_raw) if header_font_size_raw else (int(font_size_raw) if font_size_raw else 10),
                group_id=group_id,
            ))

        elif native_type in ("line_chart", "bar_chart", "column_chart"):
            categories_raw = g.get("data-native-categories", "")
            series_raw = g.get("data-native-series", "[]")
            axis_min_raw = g.get("data-native-axis-min")
            axis_max_raw = g.get("data-native-axis-max")
            axis_major_unit_raw = g.get("data-native-axis-major-unit")
            value_number_format = g.get("data-native-number-format")
            category_font_size_raw = g.get("data-native-category-font-size")
            value_font_size_raw = g.get("data-native-value-font-size")
            legend_font_size_raw = g.get("data-native-legend-font-size")

            categories = [c.strip() for c in categories_raw.split(",") if c.strip()]
            try:
                series_list = json.loads(series_raw)
            except json.JSONDecodeError as e:
                print(f"  Warning: JSON parse error in {svg_path.name} group '{group_id}': {e}")
                continue

            series_specs = []
            for s in series_list:
                series_specs.append(SeriesSpec(
                    name=s.get("name", "Series"),
                    values=[float(v) for v in s.get("values", [])],
                    color=s.get("color"),
                ))

            specs.append(ChartSpec(
                slide_index=slide_index,
                native_type=native_type,
                x=bx, y=by, width=bw, height=bh,
                categories=categories, series=series_specs,
                axis_min=float(axis_min_raw) if axis_min_raw is not None else None,
                axis_max=float(axis_max_raw) if axis_max_raw is not None else None,
                axis_major_unit=float(axis_major_unit_raw) if axis_major_unit_raw is not None else None,
                value_number_format=value_number_format,
                category_font_size=int(category_font_size_raw) if category_font_size_raw else None,
                value_font_size=int(value_font_size_raw) if value_font_size_raw else None,
                legend_font_size=int(legend_font_size_raw) if legend_font_size_raw else None,
                group_id=group_id,
            ))
        else:
            print(f"  Warning: Unknown data-native-type '{native_type}' in {svg_path.name}")

    return specs


# ── Shape removal ────────────────────────────────────────────────────────────

def _shape_in_bounds(shape, left_emu: int, top_emu: int, right_emu: int, bottom_emu: int, tolerance: int = 50000) -> bool:
    """Check if a shape's position falls within the given EMU bounds (with tolerance)."""
    try:
        sl = shape.left
        st = shape.top
        if sl is None or st is None:
            return False
        return (sl >= left_emu - tolerance and st >= top_emu - tolerance
                and sl <= right_emu + tolerance and st <= bottom_emu + tolerance)
    except Exception:
        return False


def remove_shapes_in_region(slide, x_px: float, y_px: float, w_px: float, h_px: float) -> int:
    """Remove shapes whose bounding box overlaps the given pixel region.

    Edges that only touch are treated as non-overlapping so tightly stacked
    chart/table regions do not delete each other.
    """
    left = px(x_px)
    top = px(y_px)
    right = px(x_px + w_px)
    bottom = px(y_px + h_px)

    sp_tree = slide.shapes._spTree
    to_remove = []
    for shape_elem in list(sp_tree):
        tag = shape_elem.tag.split("}")[-1] if "}" in shape_elem.tag else shape_elem.tag
        if tag not in ("sp", "grpSp", "pic", "graphicFrame", "cxnSp"):
            continue
        # Check position from XML - get bounding box
        xfrm = shape_elem.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}off")
        ext = shape_elem.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}ext")
        if xfrm is not None:
            sx = int(xfrm.get("x", 0))
            sy = int(xfrm.get("y", 0))
            ex = int(ext.get("cx", 0)) if ext is not None else 0
            ey = int(ext.get("cy", 0)) if ext is not None else 0
            shape_right = sx + ex
            shape_bottom = sy + ey
            overlaps = not (
                shape_right <= left or
                shape_bottom <= top or
                sx >= right or
                sy >= bottom
            )
            if overlaps:
                to_remove.append(shape_elem)

    for elem in to_remove:
        sp_tree.remove(elem)

    return len(to_remove)


# ── Native chart generators ─────────────────────────────────────────────────

def _parse_hex_color(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _parse_alignment(value: str):
    v = (value or "center").strip().lower()
    if v in ("left", "start"):
        return PP_ALIGN.LEFT
    if v in ("right", "end"):
        return PP_ALIGN.RIGHT
    return PP_ALIGN.CENTER


def add_native_line_chart(slide, spec: ChartSpec) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = spec.categories

    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        px(spec.x), px(spec.y), px(spec.width), px(spec.height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = len(spec.series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(spec.legend_font_size or 9)

    try:
        chart.value_axis.minimum_scale = spec.axis_min if spec.axis_min is not None else 0
        if spec.axis_max is not None:
            chart.value_axis.maximum_scale = spec.axis_max
        else:
            chart.value_axis.maximum_scale = 10000
        if spec.axis_major_unit is not None:
            chart.value_axis.major_unit = spec.axis_major_unit
        else:
            chart.value_axis.major_unit = 2000
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NEXT_TO_AXIS
        chart.value_axis.tick_labels.font.size = Pt(spec.value_font_size or 11)
        chart.value_axis.tick_labels.font.name = "Arial"
        if spec.value_number_format:
            chart.value_axis.tick_labels.number_format = spec.value_number_format
    except Exception:
        pass

    try:
        chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.NEXT_TO_AXIS
        chart.category_axis.tick_labels.font.size = Pt(spec.category_font_size or 11)
        chart.category_axis.tick_labels.font.name = "Arial"
    except Exception:
        pass

    for i, s in enumerate(spec.series):
        series = chart.series[i]
        color = _parse_hex_color(s.color) if s.color else RGBColor(0xC0, 0x00, 0x00)
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.0)
        series.marker.style = 8  # CIRCLE
        series.marker.size = 6
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = color
        series.marker.format.line.color.rgb = color
        series.has_data_labels = False


def add_native_bar_chart(slide, spec: ChartSpec) -> None:
    chart_data = CategoryChartData()
    chart_data.categories = spec.categories

    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    chart_type = XL_CHART_TYPE.BAR_STACKED if spec.native_type == "bar_chart" else XL_CHART_TYPE.COLUMN_CLUSTERED

    chart_frame = slide.shapes.add_chart(
        chart_type,
        px(spec.x), px(spec.y), px(spec.width), px(spec.height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = len(spec.series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)

    try:
        chart.plots[0].gap_width = 55
        if spec.native_type == "bar_chart":
            chart.plots[0].overlap = 100
    except Exception:
        pass

    try:
        if spec.native_type == "bar_chart":
            chart.category_axis.reverse_order = True
            chart.category_axis.tick_labels.font.size = Pt(10)
            chart.category_axis.tick_labels.font.name = "Arial"
            chart.value_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE
            chart.value_axis.has_major_gridlines = True
        else:
            chart.category_axis.tick_labels.font.size = Pt(10)
            chart.category_axis.tick_labels.font.name = "Arial"
    except Exception:
        pass

    is_percent = all(
        all(0 <= v <= 1 for v in s.values)
        for s in spec.series
    ) if spec.series else False

    for i, s in enumerate(spec.series):
        series = chart.series[i]
        if s.color:
            color = _parse_hex_color(s.color)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        series.has_data_labels = False


def add_native_pie_chart(slide, spec: ChartSpec) -> None:
    """创建环形图（doughnut），适合占比/份额类数据"""
    chart_data = ChartData()
    chart_data.categories = spec.categories

    values = spec.series[0].values if spec.series else []
    chart_data.add_series(spec.series[0].name if spec.series else "数值", values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        px(spec.x), px(spec.y), px(spec.width), px(spec.height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(spec.legend_font_size or 8)

    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '0.0%'
        data_labels.font.size = Pt(8)
    except Exception:
        pass

    all_colors = []
    if spec.series and spec.series[0].color:
        color_str = spec.series[0].color
        try:
            parsed = json.loads(color_str)
            if isinstance(parsed, list):
                all_colors = parsed
        except (json.JSONDecodeError, TypeError):
            all_colors = [color_str]

    if not all_colors:
        all_colors = ["#2B4570", "#4682B4", "#89C2D9", "#F0A500", "#E85D04",
                       "#6A994E", "#386641", "#BC4749", "#607D8B", "#8BC34A"]

    try:
        plot = chart.plots[0]
        for i in range(len(spec.categories)):
            point = plot.series[0].points[i]
            color = _parse_hex_color(all_colors[i % len(all_colors)])
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
    except Exception:
        pass


def add_native_horizontal_bar_chart(slide, spec: ChartSpec) -> None:
    """创建横向条形图（BAR_CLUSTERED），适合排名/TOP数据"""
    chart_data = CategoryChartData()
    chart_data.categories = spec.categories

    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        px(spec.x), px(spec.y), px(spec.width), px(spec.height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = len(spec.series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)

    try:
        chart.plots[0].gap_width = 80
        chart.category_axis.reverse_order = True
        chart.category_axis.tick_labels.font.size = Pt(spec.category_font_size or 9)
        chart.category_axis.tick_labels.font.name = "Microsoft YaHei"
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(spec.value_font_size or 9)
    except Exception:
        pass

    for i, s in enumerate(spec.series):
        series = chart.series[i]
        if s.color:
            color = _parse_hex_color(s.color)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        series.has_data_labels = False


def add_native_stacked_bar_chart(slide, spec: ChartSpec) -> None:
    """创建堆叠柱状图（COLUMN_STACKED），适合展示构成/占比随类别变化"""
    chart_data = CategoryChartData()
    chart_data.categories = spec.categories

    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        px(spec.x), px(spec.y), px(spec.width), px(spec.height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(spec.legend_font_size or 8)

    try:
        chart.plots[0].gap_width = 60
        chart.plots[0].overlap = 100
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(spec.value_font_size or 9)
        chart.category_axis.tick_labels.font.size = Pt(spec.category_font_size or 9)
        chart.category_axis.tick_labels.font.name = "Microsoft YaHei"
    except Exception:
        pass

    for i, s in enumerate(spec.series):
        series = chart.series[i]
        if s.color:
            color = _parse_hex_color(s.color)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
        series.has_data_labels = False


def add_native_area_chart(slide, spec: ChartSpec) -> None:
    """创建面积图（AREA），适合趋势+总量强调"""
    chart_data = CategoryChartData()
    chart_data.categories = spec.categories

    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    chart_type = XL_CHART_TYPE.AREA_STACKED if len(spec.series) > 1 else XL_CHART_TYPE.AREA

    chart_frame = slide.shapes.add_chart(
        chart_type,
        px(spec.x), px(spec.y), px(spec.width), px(spec.height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = len(spec.series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(spec.legend_font_size or 8)

    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(spec.value_font_size or 9)
        chart.category_axis.tick_labels.font.size = Pt(spec.category_font_size or 9)
        chart.category_axis.tick_labels.font.name = "Arial"
        if spec.axis_min is not None:
            chart.value_axis.minimum_scale = spec.axis_min
        if spec.axis_max is not None:
            chart.value_axis.maximum_scale = spec.axis_max
        if spec.axis_major_unit is not None:
            chart.value_axis.major_unit = spec.axis_major_unit
    except Exception:
        pass

    for i, s in enumerate(spec.series):
        series = chart.series[i]
        if s.color:
            color = _parse_hex_color(s.color)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            series.format.line.color.rgb = color


def add_native_radar_chart(slide, spec: ChartSpec) -> None:
    """创建雷达图（RADAR_FILLED），适合多维度评分对比"""
    chart_data = CategoryChartData()
    chart_data.categories = spec.categories

    for s in spec.series:
        chart_data.add_series(s.name, s.values)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR_FILLED,
        px(spec.x), px(spec.y), px(spec.width), px(spec.height),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(spec.legend_font_size or 8)

    for i, s in enumerate(spec.series):
        series = chart.series[i]
        if s.color:
            color = _parse_hex_color(s.color)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            series.format.line.color.rgb = color
            series.format.line.width = Pt(1.5)


def add_native_table(slide, spec: TableSpec) -> None:
    num_rows = len(spec.rows) + 1  # +1 header
    num_cols = len(spec.headers)
    if num_cols == 0:
        return

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        px(spec.x), px(spec.y),
        px(spec.width), px(spec.height),
    )
    table = table_shape.table

    # Column widths
    if spec.col_widths and len(spec.col_widths) == num_cols:
        for i, w in enumerate(spec.col_widths):
            table.columns[i].width = px(w)
    else:
        col_w = spec.width / num_cols
        for i in range(num_cols):
            table.columns[i].width = px(col_w)

    header_rgb = _parse_hex_color(spec.header_color)
    header_text_rgb = _parse_hex_color(spec.header_text_color)
    body_text_rgb = _parse_hex_color(spec.body_text_color)

    # Header row
    for i, header in enumerate(spec.headers):
        cell = table.cell(0, i)
        _set_cell(
            cell,
            header,
            bold=True,
            font_color=header_text_rgb,
            font_size=spec.header_font_size,
            alignment=_parse_alignment(spec.col_alignments[i] if spec.col_alignments and i < len(spec.col_alignments) else 'center'),
        )
        _fill_cell(cell, header_rgb)

    # Data rows
    for row_idx, row_data in enumerate(spec.rows):
        bg_hex = spec.row_colors[row_idx % len(spec.row_colors)] if spec.row_colors else "#FFFFFF"
        bg_rgb = _parse_hex_color(bg_hex)
        for col_idx in range(min(len(row_data), num_cols)):
            cell = table.cell(row_idx + 1, col_idx)
            _set_cell(
                cell,
                row_data[col_idx],
                font_color=body_text_rgb,
                font_size=spec.font_size,
                alignment=_parse_alignment(spec.col_alignments[col_idx] if spec.col_alignments and col_idx < len(spec.col_alignments) else 'center'),
            )
            _fill_cell(cell, bg_rgb)


def _set_cell(cell, text: str, bold: bool = False, font_color: RGBColor | None = None,
              font_size: int = 10, alignment=PP_ALIGN.CENTER) -> None:
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.alignment = alignment
        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = "Microsoft YaHei"
            if font_color:
                run.font.color.rgb = font_color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _fill_cell(cell, rgb: RGBColor) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb


# ── Main pipeline ────────────────────────────────────────────────────────────

def find_latest_pptx(exports_dir: Path) -> Path | None:
    """Find the most recently created native PPTX (excludes _svg.pptx)."""
    candidates = [
        f for f in sorted(exports_dir.glob("*.pptx"), key=lambda p: p.stat().st_mtime, reverse=True)
        if not f.name.startswith("~$") and "_svg" not in f.stem
    ]
    return candidates[0] if candidates else None


def _copy_master_from_template(pptx_path: Path, template_path: Path, verbose: bool = False) -> None:
    """Copy slide master (with logo) from template PPTX to target PPTX at ZIP level."""
    temp_dir = Path(tempfile.mkdtemp())
    try:
        # Extract target PPTX
        with zipfile.ZipFile(pptx_path, 'r') as zf:
            zf.extractall(temp_dir / 'target')

        # Extract template PPTX
        with zipfile.ZipFile(template_path, 'r') as zf:
            zf.extractall(temp_dir / 'template')

        # Copy slideMasters directory from template
        template_masters = temp_dir / 'template' / 'ppt' / 'slideMasters'
        target_masters = temp_dir / 'target' / 'ppt' / 'slideMasters'

        if template_masters.exists():
            # Copy all master XML files
            for master_file in template_masters.glob('*.xml'):
                shutil.copy2(master_file, target_masters / master_file.name)
                if verbose:
                    print(f"    Copied master: {master_file.name}")

            # Copy master relationship file
            template_masters_rels = temp_dir / 'template' / 'ppt' / 'slideMasters' / '_rels'
            target_masters_rels = temp_dir / 'target' / 'ppt' / 'slideMasters' / '_rels'
            if template_masters_rels.exists():
                target_masters_rels.mkdir(parents=True, exist_ok=True)
                for rel_file in template_masters_rels.glob('*.rels'):
                    shutil.copy2(rel_file, target_masters_rels / rel_file.name)

            # Copy master media (logos, etc.)
            template_media = temp_dir / 'template' / 'ppt' / 'media'
            target_media = temp_dir / 'target' / 'ppt' / 'media'
            if template_media.exists():
                target_media.mkdir(parents=True, exist_ok=True)
                for media_file in template_media.glob('*'):
                    if not (target_media / media_file.name).exists():
                        shutil.copy2(media_file, target_media / media_file.name)
                        if verbose:
                            print(f"    Copied media: {media_file.name}")

            # Update [Content_Types].xml to include master parts
            content_types_path = temp_dir / 'target' / '[Content_Types].xml'
            if content_types_path.exists():
                with open(content_types_path, 'r', encoding='utf-8') as f:
                    content_types_xml = f.read()

                # Add slideMaster content type if not present
                if 'slideMaster' not in content_types_xml:
                    content_types_xml = content_types_xml.replace(
                        '</Types>',
                        '  <Default ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml" Extension="xml"/>\n</Types>'
                    )
                    with open(content_types_path, 'w', encoding='utf-8') as f:
                        f.write(content_types_xml)

        # Repackage PPTX
        with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in (temp_dir / 'target').rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(temp_dir / 'target')
                    zf.write(file_path, arcname)

        if verbose:
            print(f"  Master copied from template")

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def process_project(project_path: Path, pptx_path: Path | None = None, template_path: Path | None = None, verbose: bool = True) -> bool:
    """Main entry: scan SVGs for data-native-* markers, patch the PPTX."""
    svg_dir = project_path / "svg_final"
    if not svg_dir.exists():
        svg_dir = project_path / "svg_output"
    if not svg_dir.exists():
        print(f"Error: No SVG directory found in {project_path}")
        return False

    svg_files = sorted(svg_dir.glob("*.svg"))
    if not svg_files:
        print("Error: No SVG files found")
        return False

    # Collect all native specs
    all_specs: list[ChartSpec | TableSpec] = []
    for i, svg_file in enumerate(svg_files):
        specs = extract_native_specs(svg_file, slide_index=i)
        all_specs.extend(specs)

    if not all_specs:
        if verbose:
            print("[native_charts] No data-native-* markers found in SVG files, skipping.")
        return True

    # Find PPTX
    if pptx_path is None:
        exports_dir = project_path / "exports"
        pptx_path = find_latest_pptx(exports_dir) if exports_dir.exists() else None

    if pptx_path is None or not pptx_path.exists():
        print(f"Error: No PPTX file found in exports/")
        return False

    if verbose:
        print(f"[native_charts] Processing: {pptx_path.name}")
        print(f"  Found {len(all_specs)} native element(s) across {len(svg_files)} SVG file(s)")

    # Copy master from template if provided (low-level ZIP operation)
    if template_path and template_path.exists():
        if verbose:
            print(f"  Copying master from template: {template_path.name}")
        _copy_master_from_template(pptx_path, template_path, verbose)

    prs = Presentation(str(pptx_path))

    charts_added = 0
    tables_added = 0

    for spec in all_specs:
        if spec.slide_index >= len(prs.slides):
            print(f"  Warning: slide index {spec.slide_index} out of range, skipping")
            continue

        slide = prs.slides[spec.slide_index]

        # Remove static shapes in the region
        removed = remove_shapes_in_region(slide, spec.x, spec.y, spec.width, spec.height)

        if isinstance(spec, ChartSpec):
            if spec.native_type == "line_chart":
                add_native_line_chart(slide, spec)
                charts_added += 1
                if verbose:
                    print(f"  Slide {spec.slide_index + 1}: line_chart '{spec.group_id}' (removed {removed} shapes)")
            elif spec.native_type in ("bar_chart", "column_chart"):
                add_native_bar_chart(slide, spec)
                charts_added += 1
                if verbose:
                    print(f"  Slide {spec.slide_index + 1}: bar_chart '{spec.group_id}' (removed {removed} shapes)")
        elif isinstance(spec, TableSpec):
            add_native_table(slide, spec)
            tables_added += 1
            if verbose:
                print(f"  Slide {spec.slide_index + 1}: table '{spec.group_id}' (removed {removed} shapes)")

    prs.save(str(pptx_path))

    if verbose:
        print(f"\n[Done] Patched: {pptx_path}")
        print(f"  Charts: {charts_added}, Tables: {tables_added}")

    return True


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Post-process PPTX: replace SVG-based charts/tables with native PowerPoint objects.",
    )
    parser.add_argument("project_path", help="Project directory path")
    parser.add_argument("--pptx", type=str, default=None, help="Explicit PPTX path to patch")
    parser.add_argument("--template", type=str, default=None, help="Template PPTX path (to copy master/logo)")
    parser.add_argument("-q", "--quiet", action="store_true", help="Quiet mode")

    args = parser.parse_args()
    project_path = Path(args.project_path)

    if not project_path.exists():
        print(f"Error: Path does not exist: {project_path}")
        sys.exit(1)

    pptx_path = Path(args.pptx) if args.pptx else None
    template_path = Path(args.template) if args.template else None
    ok = process_project(project_path, pptx_path=pptx_path, template_path=template_path, verbose=not args.quiet)
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    main()
